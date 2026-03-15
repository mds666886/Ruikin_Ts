from pathlib import Path


def _build_slide(page_no: int, text: str):
    text = (text or "").strip()
    title = text.splitlines()[0][:30] if text else f"第{page_no}页"
    text_count = len(text)
    risk_tags = []
    if text_count > 260:
        risk_tags.append("文字过密")
    if not text:
        risk_tags.append("内容过少")
    return {
        "page_no": page_no,
        "title_guess": title or f"第{page_no}页",
        "text_count": text_count,
        "image_ratio": 0.3,
        "is_key_page": 1 if page_no in (1, 2) else 0,
        "risk_tags": risk_tags,
        "logic_score": max(30.0, 100.0 - len(risk_tags) * 25),
        "summary": text[:120],
    }


def parse_txt(path: str):
    content = Path(path).read_text(encoding="utf-8", errors="ignore")
    chunks = [c.strip() for c in content.split("\n\n") if c.strip()]
    if not chunks:
        chunks = [""]
    return [_build_slide(i + 1, c) for i, c in enumerate(chunks)]


def parse_pdf(path: str):
    try:
        import fitz
    except Exception:
        return [_build_slide(1, "PDF 解析依赖未安装，请先安装 PyMuPDF")]

    doc = fitz.open(path)
    slides = []
    for i, page in enumerate(doc):
        slides.append(_build_slide(i + 1, page.get_text("text")))
    return slides or [_build_slide(1, "")]


def parse_pptx(path: str):
    try:
        from pptx import Presentation
    except Exception:
        return [_build_slide(1, "PPT 解析依赖未安装，请先安装 python-pptx")]

    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        slides.append(_build_slide(i + 1, "\n".join(parts)))
    return slides or [_build_slide(1, "")]


def parse_file(path: str, file_type: str):
    if file_type == "txt":
        return parse_txt(path)
    if file_type == "pdf":
        return parse_pdf(path)
    if file_type == "pptx":
        return parse_pptx(path)
    return [_build_slide(1, "暂不支持的文件类型")]
